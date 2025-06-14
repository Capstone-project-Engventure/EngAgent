import os
import re
import uuid
import json
from pathlib import Path
from datetime import datetime

import pdfplumber
import docx2txt
from pptx import Presentation
import pandas as pd

from langchain.embeddings import VertexAIEmbeddings, HuggingFaceEmbeddings
from langchain.vectorstores import Chroma
from langchain.schema import Document


# -------------------------
# Configuration
# -------------------------
DATA_DIR = Path("data")
CRAWL_DATE = datetime.utcnow().isoformat()

# Source directories and files
SOURCES = {
    "grammar": {
        "path": DATA_DIR / "bank_exercises/grammar",
        "exts": [".pdf", ".docx"],
        "type": "grammar",
    },
    "vocab_csv": {
        "path": DATA_DIR / "vocabs" / "vocab.csv",
        "type": "vocabulary",
    },
    "reading_slides": {
        "path": DATA_DIR / "Reading_Comprehension_Slides",
        "exts": [".pptx"],
        "type": "reading",
    },
    "listening_transcripts": {
        "path": DATA_DIR / "listening" / "Listening_Transcripts",
        "exts": [".txt"],
        "type": "listening",
    },
    "writing_exams": {
        "path": DATA_DIR  / "exam_bank",
        "exts": [".pdf", ".docx"],
        "type": "writing",
    },
    "speaking_scripts": {
        "path": DATA_DIR  / "Speaking_Scripts.md",
        "type": "speaking",
    },
}


# -------------------------
# Extraction functions
# -------------------------
def extract_pdf(path: Path) -> str:
    """Extract plain text from a PDF file."""
    text = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text

def extract_two_column_pdf(path: Path) -> str:
    text = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            left = page.crop((0, 0, page.width/2, page.height)).extract_text() or ""
            right = page.crop((page.width/2, 0, page.width, page.height)).extract_text() or ""
            # Nối theo đúng thứ tự xuất bản: left dòng 1, right dòng 1, left dòng 2, ...
            left_lines  = left.splitlines()
            right_lines = right.splitlines()
            for l,r in zip(left_lines, right_lines):
                text += l + " " + r + "\n"
    return text


def extract_docx(path: Path) -> str:
    """Extract plain text from a DOCX file."""
    return docx2txt.process(str(path))


def extract_pptx(path: Path) -> str:
    """Extract concatenated text from PPTX slides."""
    text = ""
    prs = Presentation(str(path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


# -------------------------
# Preprocessing & Chunking
# -------------------------

def detect_pdf_layout(path: Path) -> str:
    # Mở trang đầu, xem số vùng text box
    with pdfplumber.open(path) as pdf:
        page = pdf.pages[0]
        # nếu crop nửa trái + phải đều có text → 2 cột
        left = page.crop((0, 0, page.width/2, page.height)).extract_text() or ""
        right = page.crop((page.width/2, 0, page.width, page.height)).extract_text() or ""
        return "two-column" if left.strip() and right.strip() else "one-column"

def clean_text(txt: str) -> str:
    """Remove headers, footers, ads and normalize unicode."""
    lines = [l.strip() for l in txt.splitlines()]
    lines = [l for l in lines if len(l.split()) > 2]
    text = "\n".join(lines)
    return text.replace("ADVERTISEMENT", "")


def chunk_text(txt: str, min_words=100, max_words=300, overlap=0.2) -> list[str]:
    """Split text into overlapping chunks preserving context."""
    words = txt.split()
    size = max_words
    step = int(size * (1 - overlap))
    chunks = []
    for i in range(0, len(words), step):
        piece = words[i : i + size]
        if len(piece) >= min_words:
            chunks.append(" ".join(piece))
    return chunks
def parse_grammar_exercise(path: Path, layout: str) -> list[dict]:
    """
    Parse a grammar exercise PDF into Q&A chunks using regex.
    layout: "one-column" hoặc "two-column"
    """
    # 1) chọn extractor theo layout
    if layout == "two-column":
        raw = extract_two_column_pdf(path)
    else:
        raw = extract_pdf(path)

    # 2) tách phần Q&A và Answer Key
    parts = re.split(r"Answer Key:", raw, flags=re.IGNORECASE)
    if len(parts) < 2:
        return []
    qa_text, key_text = parts[0], parts[1]

    # 3) build map số câu => đáp án
    key_map = {
        m.group(1): m.group(2)
        for m in re.finditer(r"(\d+):\s*([A-D])", key_text)
    }

    # 4) regex lấy từng block Q&A
    pattern = re.compile(
        r"(\d{1,2})[\.\)]\s*"
        r"(.*?)\s*A\)\s*(.*?)\s*"
        r"B\)\s*(.*?)\s*"
        r"C\)\s*(.*?)\s*"
        r"D\)\s*(.*?)(?=\n\d+\.|$)",
        re.S
    )

    # 5) infer level từ tên file
    level_match = re.search(r"([a-z]\d-[a-z]\d)", path.name, re.I)
    level = level_match.group(1).upper() if level_match else None

    chunks: list[dict] = []
    for num, q, a, b, c, d in pattern.findall(qa_text):
        opts = [a.strip(), b.strip(), c.strip(), d.strip()]
        letter = key_map.get(num)
        idx = ord(letter) - ord('A') if letter else None
        answer = opts[idx] if idx is not None and 0 <= idx < len(opts) else None

        # 6) build field text để index: question + options + answer
        text = "\n".join([
            f"{num}. {q.strip()}",
            *(f"{k}) {v}" for k, v in zip(["A","B","C","D"], opts)),
            f"Answer: {answer or ''}"
        ])

        chunks.append({
            "id": str(uuid.uuid4()),
            "source": "grammar",
            "url": f"file://{path.relative_to(DATA_DIR)}",
            "crawl_date": CRAWL_DATE,
            "type": "grammar",
            "name": path.stem,
            "level": level,
            "question": q.strip(),
            "options": {"A": opts[0], "B": opts[1], "C": opts[2], "D": opts[3]},
            "answer": answer,
            "text": text,   # ← trường dùng để embedding / indexing
        })

    return chunks
# -------------------------
# Pipeline stages
# -------------------------

def collect_and_process() -> list[dict]:
    all_chunks: list[dict] = []
    cfg = SOURCES.get("grammar")
    if not cfg:
        return all_chunks

    for ext in cfg["exts"]:
        for path in cfg["path"].glob(f"*{ext}"):
            layout = detect_pdf_layout(path)  # "one-column" hoặc "two-column"
            all_chunks.extend(parse_grammar_exercise(path, layout))

    # TODO: xử lý thêm các nguồn khác (csv, pptx, txt, v.v.)
    return all_chunks

def save_chunks(chunks: list[dict], path: Path):
    """Serialize chunk list to JSON file."""
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as f:
        json.dump(chunks, f, ensure_ascii=False, indent=2)
    print(f"Saved {len(chunks)} chunks to {path}")


def build_vector_store(chunks_file: Path, persist_dir: Path, embedding):
    """Load JSON chunks, create embeddings, and persist Chroma index."""
    with chunks_file.open("r", encoding="utf-8") as f:
        chunks = json.load(f)
    # docs = [Document(page_content=c["text"], metadata=c) for c in chunks]
    for c in chunks:
    # dùng c["text"] làm nội dung index
        docs = []
        docs.append(Document(
            page_content=c["text"],
            metadata={
                "id": c["id"],
                "source": c["source"],
                "url": c["url"],
                "type": c["type"],
                "level": c["level"],
                "name": c["name"],
                # Không cần options+answer lặp lại trong metadata nếu đã có trong text
            }
        ))
    # embedder = HuggingFaceEmbeddings() # Giải thích model_name="sentence-transformers/all-MiniLM-L6-v2"
    vectordb = Chroma.from_documents(
        docs, embedding=embedding, persist_directory=str(persist_dir)
    )
    vectordb.persist()
    print("Vector store built & persisted.")
    return vectordb
 


# -------------------------
# Main pipeline entrypoint
# -------------------------
# if __name__ == "__main__":
#     chunks = collect_and_process()
#     save_chunks(chunks, DATA_DIR / "chunks.json")
#     build_vector_store(DATA_DIR / "chunks.json", DATA_DIR / "chroma_db")
